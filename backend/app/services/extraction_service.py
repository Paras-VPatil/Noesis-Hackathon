import io
from typing import Optional

import easyocr
import fitz  # PyMuPDF
import numpy as np
from PIL import Image

try:
    import docx
except ImportError:  # pragma: no cover - environment dependent
    docx = None

try:
    import pptx
except ImportError:  # pragma: no cover - environment dependent
    pptx = None

OCR_TEXT_THRESHOLD = 20
OCR_LANGUAGES = ["en"]
_ocr_reader: Optional[easyocr.Reader] = None
OCR_CONFIG_DEFAULT = {"detail": 1, "paragraph": True}
OCR_CONFIG_HANDWRITING = {
    "detail": 1,
    "paragraph": False,
    "contrast_ths": 0.05,
    "adjust_contrast": 0.7,
    "text_threshold": 0.4,
    "low_text": 0.2,
}


def _get_ocr_reader() -> easyocr.Reader:
    """Lazily initialize EasyOCR reader to avoid startup overhead."""
    global _ocr_reader
    if _ocr_reader is None:
        _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=False)
    return _ocr_reader


def _decode_text(file_bytes: bytes) -> str:
    """Decode text robustly with UTF-8 first, then latin-1 fallback."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1", errors="ignore")


def _run_easyocr(img_bytes: bytes) -> tuple[str, float]:
    """
    Run OCR on image bytes and return extracted text with average confidence.

    Returns:
        tuple[str, float]: (text, confidence)
    """
    reader = _get_ocr_reader()
    image = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    image_np = np.array(image)
    gray_np = np.array(image.convert("L"))
    high_contrast_np = np.clip((gray_np.astype(np.float32) * 1.35), 0, 255).astype(np.uint8)

    def parse_results(results: list) -> tuple[str, float]:
        if not results:
            return "", 0.0

        lines: list[str] = []
        confidences: list[float] = []

        for result in results:
            if len(result) < 3:
                continue
            text = str(result[1]).strip()
            conf = float(result[2])
            if text:
                lines.append(text)
                confidences.append(conf)

        if not lines:
            return "", 0.0

        avg_conf = sum(confidences) / len(confidences) if confidences else 0.0
        return "\n".join(lines), avg_conf

    # Pass 1: default OCR on original RGB page
    text, conf = parse_results(reader.readtext(image_np, **OCR_CONFIG_DEFAULT))
    if text:
        return text, conf

    # Pass 2: handwriting-sensitive OCR on original RGB page
    text, conf = parse_results(reader.readtext(image_np, **OCR_CONFIG_HANDWRITING))
    if text:
        return text, conf

    # Pass 3: handwriting-sensitive OCR on high-contrast grayscale
    text, conf = parse_results(reader.readtext(high_contrast_np, **OCR_CONFIG_HANDWRITING))
    if text:
        return text, conf

    return "", 0.0


async def extract_text_from_file(file_bytes: bytes, filename: str) -> dict:
    """Route file to the appropriate extraction method based on extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return await extract_pdf(file_bytes)
    if ext in ("png", "jpg", "jpeg"):
        return await extract_image(file_bytes)
    if ext == "docx":
        return extract_docx(file_bytes)
    if ext == "pptx":
        return extract_pptx(file_bytes)
    if ext == "txt":
        return {"text": _decode_text(file_bytes), "ocr_used": False, "confidence": 1.0}
    raise ValueError(f"Unsupported document format: {ext}")


async def extract_pdf(file_bytes: bytes) -> dict:
    """Extract text from PDF; use EasyOCR fallback for sparse/scanned pages."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    ocr_used = False
    ocr_confidences: list[float] = []

    for page_num, page in enumerate(doc):
        text = page.get_text().strip()

        # Sparse native text layer likely indicates scanned page.
        if len(text) < OCR_TEXT_THRESHOLD:
            ocr_used = True
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes("png")
            result = await extract_image(img_bytes, page_ref=f"p.{page_num + 1}")
            ocr_text = result.get("text", "").strip()
            ocr_confidences.append(float(result.get("confidence", 0.0)))

            # If OCR returns little/no text, preserve sparse native extraction.
            final_text = ocr_text if ocr_text else text
            pages.append({"page": page_num + 1, "text": final_text, "ocr": bool(ocr_text)})
        else:
            pages.append({"page": page_num + 1, "text": text, "ocr": False})

    full_text = "\n\n".join(
        f"[Page {p['page']}]\n{p['text']}" for p in pages if p["text"].strip()
    )
    avg_conf = (sum(ocr_confidences) / len(ocr_confidences)) if ocr_confidences else 1.0

    return {
        "text": full_text,
        "ocr_used": ocr_used,
        "pages": pages,
        "confidence": avg_conf,
    }


async def extract_image(img_bytes: bytes, page_ref: str = "img") -> dict:
    """Extract text from images using EasyOCR."""
    try:
        text, conf = _run_easyocr(img_bytes)
        return {"text": text, "ocr_used": True, "confidence": conf, "pageRef": page_ref}
    except Exception as e:
        return {
            "text": "",
            "ocr_used": True,
            "confidence": 0.0,
            "pageRef": page_ref,
            "error": str(e),
        }


def extract_docx(file_bytes: bytes) -> dict:
    """Extract text from Word Document."""
    if docx is None:
        raise ValueError("python-docx is not installed; cannot parse DOCX files.")
    doc_io = io.BytesIO(file_bytes)
    doc = docx.Document(doc_io)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return {"text": text, "ocr_used": False, "confidence": 1.0}


def extract_pptx(file_bytes: bytes) -> dict:
    """Extract text from PowerPoint slides."""
    if pptx is None:
        raise ValueError("python-pptx is not installed; cannot parse PPTX files.")
    prs_io = io.BytesIO(file_bytes)
    prs = pptx.Presentation(prs_io)

    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_content))

    full_text = "\n\n".join(slides_text)
    return {"text": full_text, "ocr_used": False, "confidence": 1.0}
